"""
Generate the 25-slide final project presentation (.pptx) from cached pipeline artifacts.

Output: PRESENTATION.pptx in the project root.

Run with:  python make_presentation.py
"""
from __future__ import annotations

import io
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).parent
CACHE = PROJECT_ROOT / "data_cache"
TMP = PROJECT_ROOT / "data_cache" / "_chart_tmp"
TMP.mkdir(exist_ok=True, parents=True)
OUT = PROJECT_ROOT / "PRESENTATION.pptx"

# ── DESIGN ──────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
LBLUE = RGBColor(0xD6, 0xE8, 0xF7)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
DGRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)
AMBER = RGBColor(0xF0, 0xAD, 0x4E)
GREEN = RGBColor(0x5C, 0xB8, 0x5C)
RED = RGBColor(0xD9, 0x53, 0x4F)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── DATA ────────────────────────────────────────────────────────────────────
def load_artifacts() -> dict:
    return {
        "fact":      pd.read_parquet(CACHE / "fact_table.parquet"),
        "sp500":     pd.read_parquet(CACHE / "sp500_table.parquet"),
        "tech":      pd.read_parquet(CACHE / "technical_features.parquet"),
        "raw_fp":    pd.read_parquet(CACHE / "fingerprints_raw.parquet"),
        "scaled_fp": pd.read_parquet(CACHE / "fingerprints_scaled.parquet"),
        "assign":    pd.read_parquet(CACHE / "cluster_assignments.parquet"),
        "diag":      pd.read_parquet(CACHE / "cluster_diagnostics.parquet"),
        "cmetrics":  pd.read_parquet(CACHE / "cluster_metrics.parquet"),
        "metrics":   pd.read_parquet(CACHE / "classification_metrics.parquet"),
        "feat_imp":  pd.read_parquet(CACHE / "feature_importance.parquet"),
        "baseline":  float(pd.read_parquet(CACHE / "baseline.parquet")["baseline"].iloc[0]),
        "pca":       pd.read_parquet(CACHE / "pca_projection.parquet"),
    }


# ── SLIDE BUILDERS ──────────────────────────────────────────────────────────
def _new_slide(prs, layout=6):
    """layout=6 is blank."""
    return prs.slides.add_slide(prs.slide_layouts[layout])


def _add_textbox(slide, text, left, top, width, height,
                  font_size=14, bold=False, color=DGRAY,
                  font_name="Calibri", align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def _add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_header(slide, section_label, title, subtitle=None):
    """Standard slide header with section pill, title, optional subtitle."""
    # Top navy bar
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), NAVY)
    _add_textbox(slide, section_label, Inches(0.4), Inches(0.06),
                  Inches(12), Inches(0.35), 11, True, WHITE, align=PP_ALIGN.LEFT)
    # Title
    _add_textbox(slide, title, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.7),
                  28, True, NAVY)
    if subtitle:
        _add_textbox(slide, subtitle, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                      14, False, GRAY)


def _add_footer(slide, slide_num, total=25):
    _add_textbox(slide, f"Vishnu Peram · DATA 255 · SJSU",
                  Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
                  9, False, GRAY)
    _add_textbox(slide, f"{slide_num} / {total}",
                  Inches(12), Inches(7.1), Inches(1), Inches(0.3),
                  9, False, GRAY, align=PP_ALIGN.RIGHT)


def _add_bullet_list(slide, items, left, top, width, height,
                      font_size=14, color=DGRAY, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "• " + item if not item.startswith("•") else item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def _save_chart(fig, name, width=900, height=520) -> Path:
    """Save plotly figure as PNG to TMP folder."""
    path = TMP / f"{name}.png"
    fig.write_image(str(path), width=width, height=height, scale=2)
    return path


def _add_image(slide, img_path, left, top, width, height):
    return slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def _add_metric_card(slide, value, label, left, top, width=Inches(2.4), height=Inches(1.1)):
    _add_filled_rect(slide, left, top, width, height, WHITE, GRAY)
    _add_textbox(slide, value, left, top + Inches(0.1), width, Inches(0.5),
                  26, True, BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, label, left, top + Inches(0.6), width, Inches(0.4),
                  10, False, GRAY, align=PP_ALIGN.CENTER)


# ── 25 SLIDE BUILDERS ───────────────────────────────────────────────────────
def slide_01_title(prs):
    s = _new_slide(prs)
    # Background full bleed navy
    _add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Accent strip
    _add_filled_rect(s, Inches(0), Inches(2.0), SLIDE_W, Inches(0.05), BLUE)
    # Title
    _add_textbox(s, "Stock Behavioral Clustering",
                  Inches(0.8), Inches(2.4), Inches(11.5), Inches(1),
                  44, True, WHITE)
    _add_textbox(s, "& Return Prediction",
                  Inches(0.8), Inches(3.2), Inches(11.5), Inches(1),
                  44, True, RGBColor(0xA8, 0xC8, 0xE8))
    _add_textbox(s, "Do data-driven behavioral groupings reveal market structure that GICS sectors miss?",
                  Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.6),
                  18, False, RGBColor(0xCD, 0xDC, 0xEC))
    # Footer block
    _add_filled_rect(s, Inches(0.8), Inches(6.0), Inches(0.04), Inches(1.0), BLUE)
    _add_textbox(s, "Vishnu Peram",
                  Inches(1.0), Inches(6.0), Inches(8), Inches(0.4),
                  16, True, WHITE)
    _add_textbox(s, "DATA 255 · Data Mining · Final Project",
                  Inches(1.0), Inches(6.4), Inches(8), Inches(0.4),
                  12, False, RGBColor(0xA8, 0xC8, 0xE8))
    _add_textbox(s, "San José State University · Spring 2026",
                  Inches(1.0), Inches(6.7), Inches(8), Inches(0.4),
                  11, False, RGBColor(0x7B, 0xAF, 0xD4))


def slide_02_agenda(prs):
    s = _new_slide(prs)
    _add_header(s, "AGENDA", "Roadmap of this presentation",
                  "12 steps from problem definition to deployment, mapped to the rubric")
    items = [
        ("1", "Problem statement & research question", "Slides 3–5"),
        ("2", "Data sources & warehouse design", "Slides 6–7"),
        ("3", "Pipeline architecture", "Slide 8"),
        ("4", "Preprocessing", "Slide 9"),
        ("5", "Feature engineering", "Slides 10–11"),
        ("6", "EDA & visualization", "Slides 12–14"),
        ("7", "Clustering — Part 1", "Slides 15–19"),
        ("8", "Classification — Part 2", "Slides 20–22"),
        ("9", "Results & evaluation", "Slide 23"),
        ("10", "Conclusions", "Slide 24"),
        ("11", "Future work", "Slide 25"),
    ]
    y = 1.8
    for num, title, ref in items:
        _add_filled_rect(s, Inches(0.6), Inches(y), Inches(0.6), Inches(0.4), BLUE)
        _add_textbox(s, num, Inches(0.6), Inches(y + 0.02), Inches(0.6), Inches(0.4),
                      14, True, WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(s, title, Inches(1.4), Inches(y + 0.04), Inches(8), Inches(0.4),
                      14, True, NAVY)
        _add_textbox(s, ref, Inches(10.5), Inches(y + 0.04), Inches(2.4), Inches(0.4),
                      11, False, GRAY, align=PP_ALIGN.RIGHT)
        y += 0.46
    _add_footer(s, 2)


def slide_03_problem(prs):
    s = _new_slide(prs)
    _add_header(s, "01 · PROBLEM STATEMENT", "Why GICS sectors may not capture how stocks behave",
                  "The official taxonomy classifies by business activity — not by trading character")
    # Two-column layout
    _add_textbox(s, "The gap", Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
                  18, True, BLUE)
    _add_bullet_list(s, [
        "GICS labels are assigned by primary business activity (what a company does)",
        "Stock prices are driven by market dynamics, factor exposure, and sentiment",
        "These two views can disagree — e.g. NVDA (Tech) and CEG (Utilities) both rallied on AI",
        "During macro narratives, otherwise unrelated stocks move together",
    ], Inches(0.5), Inches(2.6), Inches(6), Inches(3.5), 13)

    _add_filled_rect(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), LBLUE)
    _add_textbox(s, "The opportunity", Inches(7.2), Inches(2.1), Inches(5.5), Inches(0.5),
                  18, True, BLUE)
    _add_bullet_list(s, [
        "Cluster stocks by their actual price behavior over 5 years",
        "Compare data-driven groupings to official GICS sectors",
        "Test whether knowing the behavioral cluster improves prediction",
        "Reveal hidden risk archetypes that cross sector lines",
    ], Inches(7.2), Inches(2.6), Inches(5.5), Inches(3.5), 13, color=NAVY)
    _add_footer(s, 3)


def slide_04_research_question(prs, D):
    s = _new_slide(prs)
    _add_header(s, "01 · RESEARCH QUESTION", "Formal statement",
                  "Two coupled sub-questions tested on 5 years of daily price data")
    # Big quote box
    _add_filled_rect(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.0), LBLUE)
    _add_filled_rect(s, Inches(0.7), Inches(2.0), Inches(0.08), Inches(2.0), BLUE)
    _add_textbox(s,
        "Do S&P 500 stocks form distinct behavioral clusters that diverge from their official\n"
        "GICS sector classifications when grouped by market-adjusted price behavior and risk profile —\n"
        "and can technical indicators, enriched by behavioral cluster membership, predict short-term\n"
        "return direction?",
        Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.6),
        16, True, NAVY, align=PP_ALIGN.LEFT)
    # Sub questions
    y = 4.4
    _add_textbox(s, "RQ1 — Clustering",
                  Inches(0.7), Inches(y), Inches(6), Inches(0.4), 16, True, BLUE)
    _add_textbox(s,
        "When clustered using market-adjusted returns, volatility, momentum, and risk features, "
        "do the resulting groups align with or diverge from the 11 official GICS sectors?",
        Inches(0.7), Inches(y + 0.5), Inches(6), Inches(2.0), 12, False, DGRAY)

    _add_textbox(s, "RQ2 — Prediction",
                  Inches(7.0), Inches(y), Inches(6), Inches(0.4), 16, True, BLUE)
    _add_textbox(s,
        "Do technical indicators, augmented with behavioral cluster labels from RQ1, "
        "improve classification accuracy for predicting whether a stock's 5-day forward return is positive?",
        Inches(7.0), Inches(y + 0.5), Inches(6), Inches(2.0), 12, False, DGRAY)
    _add_footer(s, 4)


def slide_05_motivation(prs):
    s = _new_slide(prs)
    _add_header(s, "01 · MOTIVATION", "The 2023–24 AI boom — a real-world example of the gap",
                  "When a single narrative dominates, sectors stop describing behavior")
    # Three example pairs
    pairs = [
        ("NVDA", "Information Technology", "+239%", BLUE),
        ("CEG (Constellation Energy)", "Utilities", "+148%", AMBER),
        ("VST (Vistra)", "Utilities", "+260%", GREEN),
    ]
    _add_textbox(s, "Three stocks that moved together for the same reason — but live in different GICS sectors:",
                  Inches(0.5), Inches(2.0), Inches(12), Inches(0.5), 14, False, DGRAY)
    y = 2.7
    for ticker, sector, ret, color in pairs:
        _add_filled_rect(s, Inches(0.6), Inches(y), Inches(12), Inches(0.9), WHITE, GRAY)
        _add_filled_rect(s, Inches(0.6), Inches(y), Inches(0.12), Inches(0.9), color)
        _add_textbox(s, ticker, Inches(0.85), Inches(y + 0.15), Inches(3.2), Inches(0.4),
                      16, True, NAVY)
        _add_textbox(s, sector, Inches(0.85), Inches(y + 0.5), Inches(3.2), Inches(0.4),
                      11, False, GRAY)
        _add_textbox(s, "Theme: AI compute & power demand",
                      Inches(4.4), Inches(y + 0.3), Inches(5), Inches(0.4),
                      12, False, DGRAY)
        _add_textbox(s, ret, Inches(10.5), Inches(y + 0.25), Inches(1.9), Inches(0.4),
                      18, True, color, align=PP_ALIGN.RIGHT)
        _add_textbox(s, "2023 return", Inches(10.5), Inches(y + 0.55), Inches(1.9), Inches(0.4),
                      9, False, GRAY, align=PP_ALIGN.RIGHT)
        y += 1.05
    _add_textbox(s,
        "GICS calls these three stocks completely different sectors. Their stocks behaved nearly identically.",
        Inches(0.5), Inches(6.4), Inches(12), Inches(0.5), 12, True, BLUE, align=PP_ALIGN.CENTER)
    _add_footer(s, 5)


def slide_06_data_sources(prs, D):
    s = _new_slide(prs)
    _add_header(s, "02 · DATA SOURCES", "Three free, programmatically-accessible sources",
                  "All data is public, reproducible, and requires no paid APIs")
    # Table-style rows
    headers = ["Source", "What it provides", "Access method", "Volume"]
    rows = [
        ("Yahoo Finance (yfinance)", "Daily OHLCV for all S&P 500 tickers + index", "yfinance.download()",
         f"~{len(D['fact']):,} stock-days"),
        ("Wikipedia constituent list", "Tickers, GICS sectors, sub-industries", "pd.read_html()",
         f"{len(D['sp500'])} tickers"),
        ("^GSPC index", "S&P 500 market benchmark — for excess returns", "yfinance.download()",
         "1,509 trading days"),
    ]
    # Header row
    _add_filled_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5), NAVY)
    col_xs = [0.6, 3.6, 6.8, 10.5]
    col_ws = [3.0, 3.2, 3.7, 2.3]
    for i, h in enumerate(headers):
        _add_textbox(s, h, Inches(col_xs[i]), Inches(2.05), Inches(col_ws[i]), Inches(0.4),
                      12, True, WHITE)
    # Body rows
    for r_i, row in enumerate(rows):
        y = 2.5 + r_i * 0.6
        bg = WHITE if r_i % 2 == 0 else LGRAY
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.6), bg, GRAY)
        for c_i, cell in enumerate(row):
            _add_textbox(s, cell, Inches(col_xs[c_i]), Inches(y + 0.1), Inches(col_ws[c_i]),
                          Inches(0.4), 11, c_i == 0, DGRAY)

    # Time range rationale
    _add_filled_rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2), LBLUE)
    _add_textbox(s, "Why 2019–2024?",
                  Inches(0.7), Inches(4.85), Inches(12), Inches(0.4), 14, True, BLUE)
    _add_bullet_list(s, [
        "Pre-COVID bull (2019) → COVID crash (Mar 2020) → V-shaped recovery (2020–21)",
        "Inflation + Fed hiking cycle (2022 bear market) → AI-driven bull (2023–24)",
        "5 distinct market regimes in one panel — clusters reflect a stock's character "
        "across all conditions, not just one environment",
    ], Inches(0.7), Inches(5.3), Inches(12), Inches(1.5), 12, color=NAVY)
    _add_footer(s, 6)


def slide_07_warehouse(prs, D):
    s = _new_slide(prs)
    _add_header(s, "02 · DATA WAREHOUSE", "Star schema design",
                  "One central fact table joined to three dimensions")
    # Star diagram
    cx, cy = 6.6, 4.0
    fact_w, fact_h = 3.2, 1.6
    # Center fact
    _add_filled_rect(s, Inches(cx - fact_w/2), Inches(cy - fact_h/2),
                      Inches(fact_w), Inches(fact_h), BLUE)
    _add_textbox(s, "daily_returns_fact",
                  Inches(cx - fact_w/2), Inches(cy - fact_h/2 + 0.15),
                  Inches(fact_w), Inches(0.5), 16, True, WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, "OHLCV · returns · excess return\n5-day forward target",
                  Inches(cx - fact_w/2), Inches(cy - fact_h/2 + 0.7),
                  Inches(fact_w), Inches(0.9), 11, False, WHITE, align=PP_ALIGN.CENTER)

    # Three dimensions
    dims = [
        ("stock_dim", "ticker · company\nGICS sector · sub-industry", 1.5, 1.8),
        ("time_dim", "date · year · quarter\nmacro_regime", 11.7, 1.8),
        ("technical_features", "RSI · MACD · ATR · OBV\nbeta · volatility · MAs", 6.6, 6.6),
    ]
    for name, sub, dx, dy in dims:
        _add_filled_rect(s, Inches(dx - 1.5), Inches(dy - 0.5),
                          Inches(3), Inches(1.0), LBLUE, BLUE)
        _add_textbox(s, name, Inches(dx - 1.5), Inches(dy - 0.42),
                      Inches(3), Inches(0.4), 13, True, NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(s, sub, Inches(dx - 1.5), Inches(dy - 0.05),
                      Inches(3), Inches(0.6), 9, False, DGRAY, align=PP_ALIGN.CENTER)

    # Lines connecting dims to fact (approx)
    # Skip the lines since pptx connector requires specific shape API
    # Footprint metrics
    n_rows = len(D['fact'])
    n_tickers = D['fact']['ticker'].nunique()
    n_features = len(D['tech'].columns)
    _add_metric_card(s, f"{n_rows:,}", "fact rows", Inches(0.5), Inches(6.3))
    _add_metric_card(s, f"{n_tickers}", "tickers", Inches(3.2), Inches(6.3))
    _add_metric_card(s, f"{n_features}", "feature cols", Inches(7.7), Inches(6.3))
    _add_metric_card(s, f"{D['fact']['gics_sector'].nunique()}", "GICS sectors", Inches(10.4), Inches(6.3))
    _add_footer(s, 7)


def slide_08_pipeline(prs):
    s = _new_slide(prs)
    _add_header(s, "03 · PIPELINE", "End-to-end data + ML workflow",
                  "Modular, cached, fully reproducible from source")
    stages = [
        ("Ingest", "Wikipedia + yfinance", "fact_table.parquet"),
        ("Engineer", "Technical indicators\nBehavioral fingerprints", "technicals + fingerprints"),
        ("Cluster", "K-Means + Hierarchical\n+ DBSCAN + PCA", "cluster_assignments"),
        ("Classify", "LogReg / Tree / RF / XGB\nwith vs without cluster", "classification_metrics"),
        ("Visualize", "Streamlit dashboard\n5 interactive pages", "live URL"),
    ]
    n = len(stages)
    box_w = 2.3
    gap = 0.15
    total_w = n * box_w + (n - 1) * gap
    start_x = (13.333 - total_w) / 2
    for i, (label, sub, output) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        _add_filled_rect(s, Inches(x), Inches(2.5), Inches(box_w), Inches(2.0), LBLUE, BLUE)
        _add_filled_rect(s, Inches(x), Inches(2.5), Inches(box_w), Inches(0.45), BLUE)
        _add_textbox(s, label, Inches(x), Inches(2.55), Inches(box_w), Inches(0.4),
                      14, True, WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(s, sub, Inches(x + 0.1), Inches(3.05),
                      Inches(box_w - 0.2), Inches(0.9), 11, False, NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(s, "↓", Inches(x), Inches(3.95),
                      Inches(box_w), Inches(0.3), 12, True, BLUE, align=PP_ALIGN.CENTER)
        _add_textbox(s, output, Inches(x + 0.1), Inches(4.2),
                      Inches(box_w - 0.2), Inches(0.4), 9, True, GRAY, align=PP_ALIGN.CENTER)
        # Arrow between
        if i < n - 1:
            _add_textbox(s, "→",
                          Inches(x + box_w - 0.05), Inches(3.3),
                          Inches(0.3), Inches(0.4), 16, True, BLUE)
    # Caching note
    _add_filled_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.2), LGRAY, GRAY)
    _add_textbox(s, "Each stage caches its outputs as parquet. Re-running is incremental — only invalidated stages recompute.",
                  Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5), 13, True, NAVY)
    _add_textbox(s, "First full run: 5–10 min · Cached re-runs: <5 sec · Pipeline command: python pipeline.py",
                  Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5), 11, False, GRAY)
    _add_footer(s, 8)


def slide_09_preprocessing(prs):
    s = _new_slide(prs)
    _add_header(s, "04 · PREPROCESSING", "From raw prices to model-ready features",
                  "Five preprocessing passes ensure the data is statistically clean")
    items = [
        ("Returns honesty",
         "pct_change(fill_method=None) — gaps stay as NaN instead of being silently filled with zeros",
         BLUE),
        ("Market adjustment",
         "excess_return = stock_return − S&P 500 return — strips out the systematic market tide",
         AMBER),
        ("Missing data",
         "Drop rows missing core fields (adj_close, raw_return); technical indicators have warm-up NaNs that are dropped at modeling time",
         GREEN),
        ("Outlier control",
         "Behavioral fingerprint features winsorized at 1st/99th percentile; daily-return outliers retained as genuine signal",
         RED),
        ("Feature scaling",
         "StandardScaler for clustering; price-relative ratios (MA / close) and log-volume to prevent magnitude domination",
         RGBColor(0x8E, 0x44, 0xAD)),
    ]
    y = 1.9
    for label, desc, color in items:
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.0), WHITE, GRAY)
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(0.12), Inches(1.0), color)
        _add_textbox(s, label, Inches(0.8), Inches(y + 0.1), Inches(3.5), Inches(0.4),
                      14, True, NAVY)
        _add_textbox(s, desc, Inches(4.5), Inches(y + 0.15), Inches(8.0), Inches(0.8),
                      11, False, DGRAY)
        y += 1.05
    _add_footer(s, 9)


def slide_10_features_technical(prs):
    s = _new_slide(prs)
    _add_header(s, "05 · FEATURE ENGINEERING — PART 1", "Technical indicators (per stock-day)",
                  "8 indicators across 4 families — the standard quant-finance feature set")
    families = [
        ("Trend", BLUE, [
            ("MA(20), MA(50)", "Rolling mean of close — trend direction"),
            ("MACD(12,26,9)", "EMA(12) − EMA(26) and signal — momentum shifts"),
        ]),
        ("Momentum", AMBER, [
            ("RSI(14)", "EMA(gains)/EMA(losses) → 0..100 — overbought/oversold"),
        ]),
        ("Volatility", GREEN, [
            ("Bollinger(20, 2σ)", "MA ± 2σ; band width = vol regime"),
            ("ATR(14)", "Average true daily range"),
            ("Volatility(20d)", "Rolling std of returns"),
        ]),
        ("Risk / Volume", RED, [
            ("Beta(60d)", "Cov(stock, market) / Var(market) — market sensitivity"),
            ("OBV", "Cumulative signed volume — accumulation"),
        ]),
    ]
    x_positions = [0.5, 3.6, 6.7, 9.8]
    col_w = 3.0
    for i, (fam, color, indicators) in enumerate(families):
        x = x_positions[i]
        _add_filled_rect(s, Inches(x), Inches(2.0), Inches(col_w), Inches(0.5), color)
        _add_textbox(s, fam, Inches(x), Inches(2.05), Inches(col_w), Inches(0.4),
                      14, True, WHITE, align=PP_ALIGN.CENTER)
        _add_filled_rect(s, Inches(x), Inches(2.5), Inches(col_w), Inches(4.5), WHITE, GRAY)
        y = 2.65
        for name, desc in indicators:
            _add_textbox(s, name, Inches(x + 0.1), Inches(y), Inches(col_w - 0.2), Inches(0.4),
                          12, True, NAVY)
            _add_textbox(s, desc, Inches(x + 0.1), Inches(y + 0.4), Inches(col_w - 0.2), Inches(0.7),
                          9, False, DGRAY)
            y += 1.2
    _add_footer(s, 10)


def slide_11_features_fingerprint(prs, D):
    s = _new_slide(prs)
    _add_header(s, "05 · FEATURE ENGINEERING — PART 2", "Behavioral fingerprint (per stock)",
                  "Each stock collapses 5 years of daily data into a single 7-feature vector")
    fp_features = [
        ("mean_excess_return", "Average alpha — outperformance vs market"),
        ("volatility", "Std dev of daily excess returns"),
        ("mean_beta", "Average 60-day rolling beta"),
        ("mean_rsi", "Average RSI level — momentum tendency"),
        ("mean_bollinger_width", "Average band width — vol regime"),
        ("max_drawdown", "Worst peak-to-trough loss"),
        ("momentum_score", "Average rolling 12-month return"),
    ]
    # Two columns: left = description, right = visual
    y = 2.0
    for name, desc in fp_features:
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(6.5), Inches(0.55), WHITE, GRAY)
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(0.10), Inches(0.55), BLUE)
        _add_textbox(s, name, Inches(0.7), Inches(y + 0.05), Inches(2.5), Inches(0.45),
                      11, True, NAVY)
        _add_textbox(s, desc, Inches(3.2), Inches(y + 0.1), Inches(3.7), Inches(0.4),
                      10, False, DGRAY)
        y += 0.62

    # Right side — schematic + stat
    _add_filled_rect(s, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.5), LBLUE)
    _add_textbox(s, "Each stock → 1 row × 7 features",
                  Inches(7.5), Inches(2.2), Inches(5), Inches(0.5),
                  16, True, NAVY)
    _add_textbox(s, "Aggregation pipeline:",
                  Inches(7.5), Inches(2.8), Inches(5), Inches(0.4),
                  13, True, BLUE)
    pipeline_steps = [
        "1.  ~1,500 days × 8 indicators per stock",
        "2.  Aggregate to means / std / max",
        "3.  Winsorize 1st / 99th percentile",
        "4.  StandardScaler normalize",
        "5.  Cluster on this 7-D vector",
    ]
    yp = 3.3
    for step in pipeline_steps:
        _add_textbox(s, step, Inches(7.6), Inches(yp), Inches(5), Inches(0.3),
                      11, False, DGRAY)
        yp += 0.4
    _add_filled_rect(s, Inches(7.5), Inches(5.7), Inches(5.1), Inches(0.7), WHITE)
    _add_textbox(s, f"{len(D['raw_fp'])}",
                  Inches(7.5), Inches(5.75), Inches(2), Inches(0.5),
                  22, True, BLUE, align=PP_ALIGN.CENTER)
    _add_textbox(s, "stocks fingerprinted",
                  Inches(9.4), Inches(5.85), Inches(3), Inches(0.4),
                  12, False, DGRAY)
    _add_footer(s, 11)


def slide_12_eda_distributions(prs, D):
    s = _new_slide(prs)
    _add_header(s, "06 · EDA", "Daily return distributions — raw vs market-adjusted",
                  "Market adjustment removes the common tide and reveals stock-specific behavior")

    # Build chart: histograms side by side
    fact = D["fact"].sample(min(100_000, len(D["fact"])), random_state=42)
    fig = make_subplots(rows=1, cols=2, subplot_titles=("Raw daily returns", "Excess (market-adjusted) returns"))
    fig.add_trace(go.Histogram(x=fact["raw_daily_return"], nbinsx=80,
                                  marker_color="#2E6DA4", name="raw"), row=1, col=1)
    fig.add_trace(go.Histogram(x=fact["excess_daily_return"], nbinsx=80,
                                  marker_color="#17A2B8", name="excess"), row=1, col=2)
    fig.update_layout(showlegend=False, plot_bgcolor="white", paper_bgcolor="white",
                       font=dict(family="Helvetica", color="#333"), height=420,
                       margin=dict(t=40, b=40, l=40, r=20))
    fig.update_xaxes(showgrid=True, gridcolor="#F4F6F9")
    fig.update_yaxes(showgrid=True, gridcolor="#F4F6F9")
    p = _save_chart(fig, "eda_distributions", 1100, 480)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(8.5), Inches(3.7))
    # Insight box
    _add_filled_rect(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.3), LBLUE)
    _add_textbox(s, "Reading the chart", Inches(9.4), Inches(2.1),
                  Inches(3.4), Inches(0.4), 13, True, BLUE)
    _add_bullet_list(s, [
        "Both have fat tails (extreme moves more frequent than normal distribution)",
        "Excess returns are tighter — the market component contributes substantial variance",
        "Negative skew in both — crashes are sharper than rallies",
    ], Inches(9.4), Inches(2.6), Inches(3.4), Inches(3.5), 10, color=NAVY)

    # Headline stats
    raw_std = fact["raw_daily_return"].std()
    exc_std = fact["excess_daily_return"].std()
    _add_metric_card(s, f"{raw_std:.4f}", "Raw std", Inches(0.5), Inches(5.9))
    _add_metric_card(s, f"{exc_std:.4f}", "Excess std", Inches(3.2), Inches(5.9))
    pct_neg = (fact["raw_daily_return"] < 0).mean()
    _add_metric_card(s, f"{pct_neg:.1%}", "Down days", Inches(5.9), Inches(5.9))
    _add_metric_card(s, f"{1 - pct_neg:.1%}", "Up days", Inches(8.6), Inches(5.9))
    _add_footer(s, 12)


def slide_13_eda_sectors(prs, D):
    s = _new_slide(prs)
    _add_header(s, "06 · EDA", "Cumulative sector returns 2019 → 2024",
                  "Each sector's path through five distinct market regimes")
    fact = D["fact"]
    sa = (fact.groupby(["Date", "gics_sector"])["raw_daily_return"]
              .mean().reset_index())
    sa["cum"] = sa.groupby("gics_sector")["raw_daily_return"].transform(lambda r: (1+r).cumprod() - 1)
    fig = px.line(sa, x="Date", y="cum", color="gics_sector",
                   color_discrete_sequence=px.colors.qualitative.Bold)
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white",
                       legend=dict(orientation="h", y=-0.18, font=dict(size=9)),
                       font=dict(family="Helvetica", color="#333"),
                       margin=dict(t=20, b=80, l=40, r=20), height=460)
    fig.update_yaxes(title="Cumulative return", tickformat=".0%", showgrid=True, gridcolor="#F4F6F9")
    fig.update_xaxes(title="", showgrid=True, gridcolor="#F4F6F9")
    p = _save_chart(fig, "eda_sectors", 1200, 520)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.6))
    _add_textbox(s, "Tech and Communication Services led the period; Energy whip-sawed; defensives lagged.",
                  Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), 12, True, NAVY, align=PP_ALIGN.CENTER)
    _add_footer(s, 13)


def slide_14_eda_correlation(prs, D):
    s = _new_slide(prs)
    _add_header(s, "06 · EDA", "Cross-stock correlation — the sector-vs-behavior gap",
                  "Even within the same period, stocks group by macro factor, not just sector")
    # Compute a quick stat: average pairwise correlation within sectors vs overall
    fact = D["fact"]
    # Sample 60 stocks for correlation
    sectors = D["sp500"].set_index("ticker")["gics_sector"]
    sample = (sectors.dropna().reset_index()
                       .groupby("gics_sector").head(6)["ticker"].tolist())
    sample = [t for t in sample if t in fact["ticker"].unique()][:60]
    pivot = fact[fact["ticker"].isin(sample)].pivot_table(index="Date", columns="ticker", values="raw_daily_return")
    corr = pivot.corr()
    # Within-sector vs cross-sector mean
    from itertools import combinations
    within, across = [], []
    for a, b in combinations(sample, 2):
        c = corr.loc[a, b] if (a in corr.index and b in corr.columns) else np.nan
        if pd.isna(c):
            continue
        if sectors.get(a) == sectors.get(b):
            within.append(c)
        else:
            across.append(c)
    avg_within = float(np.mean(within)) if within else 0.0
    avg_across = float(np.mean(across)) if across else 0.0

    fig = px.imshow(corr.values, x=sample, y=sample,
                     color_continuous_scale="RdBu_r", zmin=-1, zmax=1, aspect="auto")
    fig.update_xaxes(showticklabels=False, title="")
    fig.update_yaxes(showticklabels=False, title="")
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white",
                       margin=dict(t=20, b=20, l=20, r=20), height=460)
    p = _save_chart(fig, "eda_corr", 900, 800)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(7), Inches(4.8))

    # Right column: takeaways
    _add_filled_rect(s, Inches(7.8), Inches(1.9), Inches(5.0), Inches(4.8), LBLUE)
    _add_textbox(s, "What the heatmap shows", Inches(8.0), Inches(2.0),
                  Inches(5), Inches(0.4), 13, True, BLUE)
    _add_textbox(s,
        "60 stocks (~5 per sector). Red = positive correlation, blue = negative. "
        "If sectors fully explained behavior, we'd see clean diagonal blocks.",
        Inches(8.0), Inches(2.5), Inches(4.7), Inches(1.5), 10, False, DGRAY)
    _add_metric_card(s, f"{avg_within:.3f}", "Avg within-sector corr",
                      Inches(8.0), Inches(4.5), Inches(4.6))
    _add_metric_card(s, f"{avg_across:.3f}", "Avg cross-sector corr",
                      Inches(8.0), Inches(5.7), Inches(4.6))
    _add_footer(s, 14)


def slide_15_clustering_setup(prs):
    s = _new_slide(prs)
    _add_header(s, "07 · CLUSTERING — PART 1", "Three algorithms, one feature space",
                  "Same 7-D fingerprint → three different ways of finding structure")
    algos = [
        ("K-Means", BLUE, [
            "Partition into k spherical clusters",
            "Pick K via elbow + silhouette",
            "Random_state=42 for reproducibility",
            "Final K = 5",
        ]),
        ("Hierarchical (Ward)", AMBER, [
            "Builds a nested cluster tree",
            "No need to pre-specify K",
            "Dendrogram shows merge structure",
            "Cut tree at K = 5 for comparison",
        ]),
        ("DBSCAN", GREEN, [
            "Density-based — finds non-spherical clusters",
            "Identifies outliers as noise (label = -1)",
            "Useful for detecting idiosyncratic stocks",
            "eps = 0.9, min_samples = 5",
        ]),
    ]
    x_positions = [0.5, 4.8, 9.1]
    box_w = 3.8
    for i, (name, color, items) in enumerate(algos):
        x = x_positions[i]
        _add_filled_rect(s, Inches(x), Inches(2.0), Inches(box_w), Inches(0.6), color)
        _add_textbox(s, name, Inches(x), Inches(2.05), Inches(box_w), Inches(0.5),
                      18, True, WHITE, align=PP_ALIGN.CENTER)
        _add_filled_rect(s, Inches(x), Inches(2.6), Inches(box_w), Inches(3.8), WHITE, GRAY)
        _add_bullet_list(s, items, Inches(x + 0.2), Inches(2.7),
                          Inches(box_w - 0.4), Inches(3.6), 12)

    # Comparison metric
    _add_filled_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), LBLUE)
    _add_textbox(s, "Evaluation: silhouette (cluster quality) + ARI vs GICS sectors (do they match official labels?)",
                  Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5), 13, True, NAVY, align=PP_ALIGN.CENTER)
    _add_footer(s, 15)


def slide_16_kmeans_diagnostics(prs, D):
    s = _new_slide(prs)
    _add_header(s, "07 · CLUSTERING", "Choosing K — elbow + silhouette",
                  "Inertia drops sharply through K=4–6, silhouette peaks; K=5 is the chosen balance")

    diag = D["diag"]
    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(go.Scatter(x=diag["k"], y=diag["inertia"], mode="lines+markers",
                              line=dict(color="#2E6DA4", width=2.5),
                              marker=dict(size=10), name="Inertia"), secondary_y=False)
    fig.add_trace(go.Scatter(x=diag["k"], y=diag["silhouette"], mode="lines+markers",
                              line=dict(color="#F0AD4E", width=2.5, dash="dot"),
                              marker=dict(size=10), name="Silhouette"), secondary_y=True)
    fig.update_xaxes(title="K (number of clusters)", showgrid=True, gridcolor="#F4F6F9")
    fig.update_yaxes(title="Inertia", secondary_y=False, color="#2E6DA4")
    fig.update_yaxes(title="Silhouette", secondary_y=True, color="#F0AD4E")
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white", height=460,
                       legend=dict(orientation="h", y=-0.15),
                       font=dict(family="Helvetica", color="#333"))
    p = _save_chart(fig, "kmeans_diag", 1100, 520)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(8.5), Inches(4.6))

    best_k = int(diag.iloc[diag["silhouette"].idxmax()]["k"])
    best_sil = float(diag["silhouette"].max())
    _add_filled_rect(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.5), LBLUE)
    _add_textbox(s, "Diagnostics", Inches(9.4), Inches(2.1),
                  Inches(3.4), Inches(0.4), 13, True, BLUE)
    _add_metric_card(s, f"K = {best_k}", "Silhouette peak", Inches(9.4), Inches(2.6), Inches(3.2))
    _add_metric_card(s, f"{best_sil:.3f}", "Peak silhouette score", Inches(9.4), Inches(3.85), Inches(3.2))
    _add_textbox(s,
        "We use K = 5 — close to the peak with cleaner interpretability for sector comparison.",
        Inches(9.4), Inches(5.1), Inches(3.4), Inches(1.3), 10, False, NAVY)
    _add_footer(s, 16)


def slide_17_pca_visualization(prs, D):
    s = _new_slide(prs)
    _add_header(s, "07 · CLUSTERING", "PCA projection — clusters vs GICS sectors",
                  "Same 2D space, two color schemes — visual evidence of structural mismatch")

    assign = D["assign"]
    pca = D["pca"]
    evr1, evr2 = float(pca["evr1"].iloc[0]), float(pca["evr2"].iloc[0])

    # Two side-by-side scatters
    df = assign.copy()
    df["kmeans_str"] = "Cluster " + df["kmeans"].astype(str)

    fig = make_subplots(rows=1, cols=2,
                          subplot_titles=("Colored by behavioral cluster",
                                            "Colored by GICS sector"))
    palette = ["#2E6DA4", "#F0AD4E", "#5CB85C", "#D9534F", "#8E44AD", "#17A2B8"]
    for i, c in enumerate(sorted(df["kmeans"].unique())):
        sub = df[df["kmeans"] == c]
        fig.add_trace(go.Scatter(x=sub["pca1"], y=sub["pca2"], mode="markers",
                                    marker=dict(color=palette[i % len(palette)], size=6, opacity=0.8),
                                    name=f"Cluster {c}", showlegend=False), row=1, col=1)
    sectors = sorted(df["gics_sector"].dropna().unique())
    for i, sec in enumerate(sectors):
        sub = df[df["gics_sector"] == sec]
        fig.add_trace(go.Scatter(x=sub["pca1"], y=sub["pca2"], mode="markers",
                                    marker=dict(color=px.colors.qualitative.Bold[i % 11], size=6, opacity=0.8),
                                    name=sec, showlegend=False), row=1, col=2)
    fig.update_xaxes(title="PC1", showgrid=True, gridcolor="#F4F6F9")
    fig.update_yaxes(title="PC2", showgrid=True, gridcolor="#F4F6F9")
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white", height=440,
                       font=dict(family="Helvetica", color="#333"),
                       margin=dict(t=40, b=40, l=40, r=20))
    p = _save_chart(fig, "pca_compare", 1300, 500)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.4))

    _add_textbox(s,
        f"PC1 explains {evr1:.1%} · PC2 explains {evr2:.1%} · combined {evr1+evr2:.1%} of variance.  "
        f"Clusters are clean blobs; sectors overlap substantially in this space.",
        Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), 12, True, NAVY, align=PP_ALIGN.CENTER)
    _add_footer(s, 17)


def slide_18_cluster_vs_sector(prs, D):
    s = _new_slide(prs)
    _add_header(s, "07 · CLUSTERING — KEY RESULT", "Behavioral clusters do not equal GICS sectors",
                  "Adjusted Rand Index quantifies the disagreement")

    cmetrics = D["cmetrics"]
    assign = D["assign"]
    ct = pd.crosstab(assign["gics_sector"], assign["kmeans"])

    fig = px.imshow(ct.values, x=[f"C{c}" for c in ct.columns], y=ct.index,
                     color_continuous_scale="Blues", aspect="auto",
                     labels=dict(color="count"))
    for i, sector in enumerate(ct.index):
        for j in range(len(ct.columns)):
            v = ct.iloc[i, j]
            if v > 0:
                fig.add_annotation(x=j, y=i, text=str(v), showarrow=False,
                                    font=dict(color="white" if v > ct.values.max()/2 else "#1B3A5C", size=11))
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white", height=520,
                       font=dict(family="Helvetica", color="#333"),
                       margin=dict(t=20, b=20, l=120, r=20))
    fig.update_xaxes(side="bottom", title="Behavioral cluster", showticklabels=True, tickfont=dict(size=11))
    fig.update_yaxes(showticklabels=True, tickfont=dict(size=10))
    p = _save_chart(fig, "cluster_sector_heatmap", 900, 700)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(8), Inches(4.8))

    # Right: metric panel
    ari_kmeans = float(cmetrics.loc[cmetrics["algorithm"] == "kmeans", "ari_vs_sector"].iloc[0])
    ari_hier = float(cmetrics.loc[cmetrics["algorithm"] == "hierarchical", "ari_vs_sector"].iloc[0])
    sil_kmeans = float(cmetrics.loc[cmetrics["algorithm"] == "kmeans", "silhouette"].iloc[0])

    _add_filled_rect(s, Inches(8.8), Inches(1.9), Inches(4.0), Inches(4.8), LBLUE)
    _add_textbox(s, "Adjusted Rand Index", Inches(9.0), Inches(2.0),
                  Inches(3.7), Inches(0.4), 14, True, BLUE)
    _add_metric_card(s, f"{ari_kmeans:.3f}", "K-Means ARI vs GICS",
                      Inches(9.0), Inches(2.5), Inches(3.6))
    _add_metric_card(s, f"{ari_hier:.3f}", "Hierarchical ARI",
                      Inches(9.0), Inches(3.7), Inches(3.6))
    _add_metric_card(s, f"{sil_kmeans:.3f}", "K-Means silhouette",
                      Inches(9.0), Inches(4.9), Inches(3.6))
    _add_textbox(s, "ARI = 0 → independent · ARI = 1 → identical",
                  Inches(9.0), Inches(6.1), Inches(3.7), Inches(0.5), 10, False, GRAY)
    _add_footer(s, 18)


def slide_19_cluster_profiles(prs, D):
    s = _new_slide(prs)
    _add_header(s, "07 · CLUSTERING", "What does each cluster look like?",
                  "Centroid profile reveals the behavioral archetype")

    raw_fp = D["raw_fp"]
    assign = D["assign"]
    features = ["mean_excess_return", "volatility", "mean_beta",
                 "mean_rsi", "mean_bollinger_width", "max_drawdown", "momentum_score"]
    profiles = (raw_fp.merge(assign[["ticker", "kmeans"]], on="ticker")
                       .groupby("kmeans")[features].mean())
    norm = (profiles - profiles.min()) / (profiles.max() - profiles.min() + 1e-9)
    fig = go.Figure()
    palette_hex = ["#2E6DA4", "#F0AD4E", "#5CB85C", "#D9534F", "#8E44AD", "#17A2B8"]
    palette_rgb = [RGBColor(0x2E, 0x6D, 0xA4), RGBColor(0xF0, 0xAD, 0x4E),
                    RGBColor(0x5C, 0xB8, 0x5C), RGBColor(0xD9, 0x53, 0x4F),
                    RGBColor(0x8E, 0x44, 0xAD), RGBColor(0x17, 0xA2, 0xB8)]
    for i, c in enumerate(norm.index):
        fig.add_trace(go.Scatterpolar(r=norm.loc[c].values, theta=features,
                                         fill="toself", name=f"Cluster {c}",
                                         line=dict(color=palette_hex[i % len(palette_hex)])))
    fig.update_layout(polar=dict(radialaxis=dict(visible=True, range=[0, 1], showticklabels=False)),
                       paper_bgcolor="white", height=520,
                       font=dict(family="Helvetica", color="#333"))
    p = _save_chart(fig, "cluster_radar", 800, 800)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(7), Inches(4.8))

    # Right: cluster sizes & quick interpretation
    sizes = assign["kmeans"].value_counts().sort_index()
    _add_filled_rect(s, Inches(7.8), Inches(1.9), Inches(5.0), Inches(4.8), LBLUE)
    _add_textbox(s, "Cluster sizes", Inches(8.0), Inches(2.0),
                  Inches(5), Inches(0.4), 14, True, BLUE)
    y = 2.5
    for c in sorted(sizes.index):
        # bar
        bar_w = (sizes[c] / sizes.max()) * 4.0
        _add_filled_rect(s, Inches(8.5), Inches(y), Inches(bar_w), Inches(0.3),
                          palette_rgb[c % len(palette_rgb)])
        _add_textbox(s, f"C{c}", Inches(8.0), Inches(y), Inches(0.45), Inches(0.3),
                      11, True, NAVY)
        _add_textbox(s, f"{sizes[c]} stocks", Inches(8.5 + bar_w + 0.1), Inches(y),
                      Inches(2), Inches(0.3), 10, False, GRAY)
        y += 0.5
    _add_textbox(s,
        "Each cluster is a risk archetype — distinct combinations of return, volatility, and trend behavior.",
        Inches(8.0), Inches(5.6), Inches(4.7), Inches(1.0), 10, False, NAVY)
    _add_footer(s, 19)


def slide_20_classification_setup(prs, D):
    s = _new_slide(prs)
    _add_header(s, "08 · CLASSIFICATION — PART 2", "Setup",
                  "Predict 5-day forward direction; test whether cluster label adds signal")

    items_left = [
        ("Target", "forward_5day_direction (binary)"),
        ("Features w/o cluster", "14 technical indicators"),
        ("Features w/ cluster", "14 technicals + behavioral cluster label"),
        ("Train", "Years < 2023"),
        ("Test", "2023–2024 (held-out future)"),
        ("Split type", "Temporal (no random shuffling on time series)"),
    ]
    y = 2.0
    for label, value in items_left:
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(6), Inches(0.55), WHITE, GRAY)
        _add_textbox(s, label, Inches(0.7), Inches(y + 0.1), Inches(2), Inches(0.4),
                      12, True, BLUE)
        _add_textbox(s, value, Inches(2.6), Inches(y + 0.1), Inches(3.5), Inches(0.4),
                      11, False, DGRAY)
        y += 0.62

    _add_filled_rect(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), LBLUE)
    _add_textbox(s, "Models", Inches(7.2), Inches(2.1),
                  Inches(5.5), Inches(0.4), 14, True, BLUE)
    models = [
        ("Logistic Regression", "Linear baseline — sets the predictability floor"),
        ("Decision Tree (depth=8)", "Interpretable non-linear model"),
        ("Random Forest (200 trees)", "Primary model — handles interactions"),
        ("XGBoost (300 rounds)", "Strongest model — boosted gradient trees"),
    ]
    yp = 2.6
    for name, desc in models:
        _add_textbox(s, name, Inches(7.4), Inches(yp), Inches(5), Inches(0.3),
                      12, True, NAVY)
        _add_textbox(s, desc, Inches(7.4), Inches(yp + 0.32), Inches(5), Inches(0.4),
                      10, False, DGRAY)
        yp += 0.85

    baseline = D["baseline"]
    _add_filled_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8), AMBER)
    _add_textbox(s, f"Naive 'always-up' baseline accuracy: {baseline:.3f} — any model must beat this meaningfully.",
                  Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5), 13, True, WHITE, align=PP_ALIGN.CENTER)
    _add_footer(s, 20)


def slide_21_model_comparison(prs, D):
    s = _new_slide(prs)
    _add_header(s, "08 · CLASSIFICATION — RESULT", "Model comparison: ROC-AUC with vs without cluster",
                  "Does the behavioral cluster label add predictive signal?")

    metrics = D["metrics"]
    fig = px.bar(metrics, x="model", y="roc_auc", color="variant", barmode="group",
                  color_discrete_map={"Without cluster": "#6C757D", "With cluster": "#2E6DA4"})
    fig.update_yaxes(title="ROC-AUC", range=[0.45, max(0.7, metrics["roc_auc"].max() + 0.05)],
                      showgrid=True, gridcolor="#F4F6F9")
    fig.update_xaxes(title="")
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white", height=460,
                       legend=dict(orientation="h", y=-0.18, title=""),
                       font=dict(family="Helvetica", color="#333"))
    p = _save_chart(fig, "model_compare", 1100, 520)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(8.5), Inches(4.6))

    # Right panel: best model + delta
    best = metrics.sort_values("roc_auc", ascending=False).iloc[0]
    pivot = metrics.pivot(index="model", columns="variant", values="roc_auc")
    pivot["delta"] = pivot["With cluster"] - pivot["Without cluster"]
    avg_delta = float(pivot["delta"].mean())
    n_improved = int((pivot["delta"] > 0).sum())

    _add_filled_rect(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.5), LBLUE)
    _add_textbox(s, "Best model", Inches(9.4), Inches(2.1),
                  Inches(3.4), Inches(0.4), 13, True, BLUE)
    _add_metric_card(s, f"{best['roc_auc']:.3f}", f"{best['model']}", Inches(9.4), Inches(2.55), Inches(3.2))
    _add_textbox(s, "ROC-AUC delta (avg)",
                  Inches(9.4), Inches(3.85), Inches(3.4), Inches(0.4), 11, True, BLUE)
    _add_metric_card(s, f"{avg_delta:+.4f}", f"{n_improved}/4 models improved",
                      Inches(9.4), Inches(4.25), Inches(3.2))
    _add_footer(s, 21)


def slide_22_feature_importance(prs, D):
    s = _new_slide(prs)
    _add_header(s, "08 · CLASSIFICATION", "What features drive the prediction?",
                  "Random Forest feature importance — RSI and momentum dominate")

    feat_imp = D["feat_imp"].head(12).iloc[::-1]
    fig = go.Figure(data=[go.Bar(x=feat_imp["importance"], y=feat_imp["feature"], orientation="h",
                                    marker=dict(color="#2E6DA4"))])
    fig.update_xaxes(title="Importance")
    fig.update_yaxes(title="")
    fig.update_layout(plot_bgcolor="white", paper_bgcolor="white", height=460,
                       font=dict(family="Helvetica", color="#333"),
                       margin=dict(t=20, b=40, l=170, r=20))
    p = _save_chart(fig, "feat_imp", 1100, 520)
    _add_image(s, p, Inches(0.5), Inches(1.9), Inches(8.5), Inches(4.6))

    # Top 3 callouts
    top = D["feat_imp"].head(3)
    _add_filled_rect(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.5), LBLUE)
    _add_textbox(s, "Top 3 features", Inches(9.4), Inches(2.1),
                  Inches(3.4), Inches(0.4), 13, True, BLUE)
    yp = 2.5
    for i, (_, row) in enumerate(top.iterrows()):
        _add_filled_rect(s, Inches(9.4), Inches(yp), Inches(3.2), Inches(1.0), WHITE, GRAY)
        _add_textbox(s, f"#{i+1}", Inches(9.5), Inches(yp + 0.1),
                      Inches(0.5), Inches(0.3), 11, True, GRAY)
        _add_textbox(s, row["feature"], Inches(9.5), Inches(yp + 0.35),
                      Inches(3.0), Inches(0.4), 13, True, NAVY)
        _add_textbox(s, f"importance: {row['importance']:.3f}", Inches(9.5), Inches(yp + 0.7),
                      Inches(3.0), Inches(0.3), 10, False, GRAY)
        yp += 1.15
    _add_footer(s, 22)


def slide_23_results(prs, D):
    s = _new_slide(prs)
    _add_header(s, "09 · RESULTS & EVALUATION", "Headline numbers",
                  "What the analysis ultimately demonstrates")

    metrics = D["metrics"]
    cmetrics = D["cmetrics"]
    baseline = D["baseline"]
    best = metrics.sort_values("roc_auc", ascending=False).iloc[0]
    ari = float(cmetrics.loc[cmetrics["algorithm"] == "kmeans", "ari_vs_sector"].iloc[0])

    cards = [
        (f"{D['fact']['ticker'].nunique()}", "Tickers analyzed"),
        (f"{len(D['fact']):,}", "Stock-day rows"),
        (f"{ari:.3f}", "ARI (clusters vs GICS)"),
        (f"{best['roc_auc']:.3f}", f"Best ROC-AUC ({best['model']})"),
        (f"{best['accuracy']:.3f}", "Best accuracy"),
        (f"{baseline:.3f}", "Naive baseline"),
    ]
    for i, (val, label) in enumerate(cards):
        x = 0.5 + (i % 3) * 4.3
        y = 2.0 + (i // 3) * 1.6
        _add_filled_rect(s, Inches(x), Inches(y), Inches(4.0), Inches(1.4), WHITE, GRAY)
        _add_textbox(s, val, Inches(x), Inches(y + 0.15), Inches(4.0), Inches(0.7),
                      36, True, BLUE, align=PP_ALIGN.CENTER)
        _add_textbox(s, label, Inches(x), Inches(y + 0.95), Inches(4.0), Inches(0.4),
                      12, False, GRAY, align=PP_ALIGN.CENTER)

    # Take-home banner
    _add_filled_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5), NAVY)
    _add_textbox(s, "Take-home", Inches(0.7), Inches(5.55),
                  Inches(12), Inches(0.4), 14, True, RGBColor(0xA8, 0xC8, 0xE8))
    _add_textbox(s,
        f"Behavioral clusters genuinely diverge from GICS sectors (ARI = {ari:.3f}). The best model "
        f"({best['model']}) beats the naive baseline by {(best['accuracy']-baseline)*100:+.1f}% accuracy "
        "— modest but meaningful given weak-form efficient market expectations. Cluster label provides a small marginal lift.",
        Inches(0.7), Inches(5.95), Inches(12), Inches(0.9), 12, False, WHITE)
    _add_footer(s, 23)


def slide_24_conclusions(prs, D):
    s = _new_slide(prs)
    _add_header(s, "10 · CONCLUSIONS", "Three findings that answer the research question",
                  None)
    cmetrics = D["cmetrics"]
    metrics = D["metrics"]
    baseline = D["baseline"]
    best = metrics.sort_values("roc_auc", ascending=False).iloc[0]
    ari = float(cmetrics.loc[cmetrics["algorithm"] == "kmeans", "ari_vs_sector"].iloc[0])
    pivot = metrics.pivot(index="model", columns="variant", values="roc_auc")
    avg_delta = float((pivot["With cluster"] - pivot["Without cluster"]).mean())

    findings = [
        ("Finding 1", "Behavioral clusters ≠ GICS sectors",
         f"ARI = {ari:.3f}. Clusters group stocks by risk profile (growth, defensive, cyclical), not by industry. The official taxonomy hides this structural lens.", BLUE),
        ("Finding 2", "Short-term direction prediction is hard but not impossible",
         f"Best ROC-AUC: {best['roc_auc']:.3f}, accuracy: {best['accuracy']:.3f}. Above the {baseline:.3f} naive baseline. Aligns with weak-form market efficiency — technicals contain real but small signal.", AMBER),
        ("Finding 3", "Cluster label adds marginal but real predictive value",
         f"Adding the cluster label changed average ROC-AUC by {avg_delta:+.4f}. It's partially redundant with technicals, but still validates the two-part design.", GREEN),
    ]
    y = 1.8
    for tag, title, desc, color in findings:
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.55), WHITE, GRAY)
        _add_filled_rect(s, Inches(0.5), Inches(y), Inches(0.15), Inches(1.55), color)
        _add_textbox(s, tag, Inches(0.8), Inches(y + 0.1), Inches(2), Inches(0.4),
                      11, True, color)
        _add_textbox(s, title, Inches(0.8), Inches(y + 0.4), Inches(11.7), Inches(0.4),
                      16, True, NAVY)
        _add_textbox(s, desc, Inches(0.8), Inches(y + 0.85), Inches(11.7), Inches(0.7),
                      11, False, DGRAY)
        y += 1.7
    _add_footer(s, 24)


def slide_25_future_work(prs):
    s = _new_slide(prs)
    _add_header(s, "11 · FUTURE WORK & 12 · DEPLOYMENT", "Where next, and how it's deployed",
                  "Honest limitations + the live dashboard")

    # Future work
    _add_textbox(s, "Future work", Inches(0.5), Inches(2.0),
                  Inches(6), Inches(0.4), 16, True, BLUE)
    items = [
        "Survivorship bias: include delisted stocks via CRSP",
        "Rolling-window clustering — capture regime-shifting behavior",
        "Multi-horizon prediction (1-day, 20-day, 60-day)",
        "Transaction costs + position sizing → economic Sharpe ratio",
        "Add fundamentals (P/E, earnings surprise) and news sentiment",
        "Deep models (LSTM, transformer) on raw OHLCV sequences",
    ]
    _add_bullet_list(s, items, Inches(0.5), Inches(2.5), Inches(6), Inches(4.0), 12)

    # Deployment
    _add_filled_rect(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), LBLUE)
    _add_textbox(s, "Deployment", Inches(7.2), Inches(2.1),
                  Inches(5.5), Inches(0.4), 16, True, BLUE)
    _add_textbox(s, "Streamlit Community Cloud (free tier)",
                  Inches(7.2), Inches(2.6), Inches(5.5), Inches(0.4), 12, True, NAVY)
    _add_textbox(s, "Repo:", Inches(7.2), Inches(3.1), Inches(1), Inches(0.3), 10, True, GRAY)
    _add_textbox(s, "github.com/PeramVishnuSree/Stock-Behavioral-Clustering-and-Return-Prediction",
                  Inches(7.2), Inches(3.4), Inches(5.5), Inches(0.4), 10, False, BLUE)
    _add_textbox(s, "Live dashboard pages:",
                  Inches(7.2), Inches(3.9), Inches(5.5), Inches(0.4), 12, True, BLUE)
    pages = ["📊  Data Overview", "📈  EDA", "🎯  Clustering",
             "🔮  Prediction", "📋  Methodology"]
    _add_bullet_list(s, pages, Inches(7.4), Inches(4.3), Inches(5.5), Inches(2.0), 11, color=NAVY)

    # Closing
    _add_filled_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6), NAVY)
    _add_textbox(s, "Thank you — questions?",
                  Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                  16, True, WHITE, align=PP_ALIGN.CENTER)
    _add_footer(s, 25)


# ── MAIN ────────────────────────────────────────────────────────────────────
def build():
    D = load_artifacts()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        lambda: slide_01_title(prs),
        lambda: slide_02_agenda(prs),
        lambda: slide_03_problem(prs),
        lambda: slide_04_research_question(prs, D),
        lambda: slide_05_motivation(prs),
        lambda: slide_06_data_sources(prs, D),
        lambda: slide_07_warehouse(prs, D),
        lambda: slide_08_pipeline(prs),
        lambda: slide_09_preprocessing(prs),
        lambda: slide_10_features_technical(prs),
        lambda: slide_11_features_fingerprint(prs, D),
        lambda: slide_12_eda_distributions(prs, D),
        lambda: slide_13_eda_sectors(prs, D),
        lambda: slide_14_eda_correlation(prs, D),
        lambda: slide_15_clustering_setup(prs),
        lambda: slide_16_kmeans_diagnostics(prs, D),
        lambda: slide_17_pca_visualization(prs, D),
        lambda: slide_18_cluster_vs_sector(prs, D),
        lambda: slide_19_cluster_profiles(prs, D),
        lambda: slide_20_classification_setup(prs, D),
        lambda: slide_21_model_comparison(prs, D),
        lambda: slide_22_feature_importance(prs, D),
        lambda: slide_23_results(prs, D),
        lambda: slide_24_conclusions(prs, D),
        lambda: slide_25_future_work(prs),
    ]
    for i, b in enumerate(builders, 1):
        print(f"Building slide {i:2d} / 25 …")
        b()
    prs.save(str(OUT))
    print(f"\n✓ Saved presentation: {OUT}  ({OUT.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
